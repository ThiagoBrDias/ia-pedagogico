import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

class PPTService:
    def create_presentation(self, title: str, slides: list, output_dir: str) -> str:
        """Criar apresentação PowerPoint do zero"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide de título
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_placeholder = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title_placeholder.text = title
            subtitle.text = "Material Pedagógico"
            
            # Adicionar slides de conteúdo
            for slide_data in slides:
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                
                slide_title = slide.shapes.title
                slide_title.text = slide_data.get("title", "")
                
                # Adicionar conteúdo
                content = slide_data.get("content", "")
                if content:
                    content_box = slide.placeholders[1]
                    text_frame = content_box.text_frame
                    text_frame.text = content
            
            # Salvar apresentação
            output_path = os.path.join(output_dir, f"{title}.pptx")
            prs.save(output_path)
            
            return output_path
        except Exception as e:
            raise Exception(f"Erro ao criar apresentação: {str(e)}")
    
    def extract_text(self, file_path: str) -> list:
        """Extrair texto de um PowerPoint"""
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = {
                    "slide_number": i + 1,
                    "title": "",
                    "content": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape == slide.shapes.title:
                            slide_text["title"] = shape.text
                        else:
                            if shape.text.strip():
                                slide_text["content"].append(shape.text)
                
                slides_content.append(slide_text)
            
            return slides_content
        except Exception as e:
            raise Exception(f"Erro ao extrair texto do PowerPoint: {str(e)}")
    
    def add_slide(self, file_path: str, slide_title: str, slide_content: str, output_dir: str) -> str:
        """Adicionar slide a uma apresentação existente"""
        try:
            prs = Presentation(file_path)
            
            # Adicionar novo slide
            content_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(content_slide_layout)
            
            title = slide.shapes.title
            title.text = slide_title
            
            content_box = slide.placeholders[1]
            text_frame = content_box.text_frame
            text_frame.text = slide_content
            
            # Salvar apresentação
            output_path = os.path.join(output_dir, "updated.pptx")
            prs.save(output_path)
            
            return output_path
        except Exception as e:
            raise Exception(f"Erro ao adicionar slide: {str(e)}")
    
    def modify_slide(self, file_path: str, slide_number: int, new_title: str, new_content: str, output_dir: str) -> str:
        """Modificar um slide específico"""
        try:
            prs = Presentation(file_path)
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise Exception("Número de slide inválido")
            
            slide = prs.slides[slide_number - 1]
            
            # Modificar título
            if slide.shapes.title:
                slide.shapes.title.text = new_title
            
            # Modificar conteúdo
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape != slide.shapes.title:
                    shape.text = new_content
                    break
            
            # Salvar apresentação
            output_path = os.path.join(output_dir, "modified.pptx")
            prs.save(output_path)
            
            return output_path
        except Exception as e:
            raise Exception(f"Erro ao modificar slide: {str(e)}")
    
    def get_presentation_info(self, file_path: str) -> dict:
        """Obter informações da apresentação"""
        try:
            prs = Presentation(file_path)
            return {
                "num_slides": len(prs.slides),
                "width": prs.slide_width,
                "height": prs.slide_height
            }
        except Exception as e:
            raise Exception(f"Erro ao obter informações da apresentação: {str(e)}")

